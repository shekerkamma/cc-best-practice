"""
Generates ai-pitch-generator-demo.pptx — an 18-slide PowerPoint deck
matching the HTML presentation structure and dark-navy DevGTM color palette.

This is a meta-demo: a sample sales pitch that demonstrates what DevGTM Mesh's
AI Pitch Generator (Use Case #2) actually produces. Fictional prospect: Lumen Analytics.

Run: python3 build-deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys

# ── Color palette ─────────────────────────────────────────
BG       = RGBColor(0x0b, 0x0f, 0x1a)   # dark navy
CARD     = RGBColor(0x13, 0x19, 0x29)   # card bg
BORDER   = RGBColor(0x1e, 0x2d, 0x47)   # borders
TEXT     = RGBColor(0xe8, 0xed, 0xf5)   # primary text
MUTED    = RGBColor(0x6b, 0x7f, 0xa3)   # muted text
BODY     = RGBColor(0xa8, 0xbd, 0xd4)   # body text

PURPLE   = RGBColor(0x7c, 0x3a, 0xed)   # primary purple
PURPLE_L = RGBColor(0xa7, 0x8b, 0xfa)   # light purple
TEAL     = RGBColor(0x2d, 0xd4, 0xbf)   # teal accent
AMBER    = RGBColor(0xf5, 0x9e, 0x0b)   # amber
BLUE     = RGBColor(0x3b, 0x82, 0xf6)   # blue
BLUE_L   = RGBColor(0x60, 0xa5, 0xfa)   # light blue
GREEN    = RGBColor(0x34, 0xd3, 0x99)   # green
PINK     = RGBColor(0xf4, 0x72, 0xb6)   # pink

RED_L    = RGBColor(0xf8, 0x71, 0x71)   # light red (danger)
GREEN_L  = RGBColor(0x4a, 0xde, 0x80)   # light green (success)

TOTAL_SLIDES = 18

# ── Slide dimensions (widescreen 16:9) ────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_logo(slide):
    """Top-left logo: colored square + DEVGTM MESH label."""
    add_rect(slide, Inches(0.4), Inches(0.25), Inches(0.35), Inches(0.35), PURPLE)
    add_textbox(slide, Inches(0.82), Inches(0.25), Inches(2.5), Inches(0.35),
                "DEVGTM MESH", size=11, bold=True, color=PURPLE_L)


def slide_number(slide, n, total=TOTAL_SLIDES):
    add_textbox(slide, Inches(11.8), Inches(7.05), Inches(1.4), Inches(0.35),
                f"{n} / {total}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def section_label(slide, label, color=PURPLE_L, left=Inches(0.4), top=Inches(0.75)):
    add_textbox(slide, left, top, Inches(5), Inches(0.35),
                label.upper(), size=11, bold=True, color=color)


def ai_badge(slide):
    add_rect(slide, Inches(4.5), Inches(0.2), Inches(4.3), Inches(0.45), CARD, PURPLE)
    add_textbox(slide, Inches(4.5), Inches(0.22), Inches(4.3), Inches(0.41),
                "✦  AI-GENERATED PITCH  ·  PITCH GENERATOR DEMO",
                size=11, bold=True, color=PURPLE_L, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────
# SLIDE 01 — Cover
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
ai_badge(sl)

add_textbox(sl, Inches(1.5), Inches(1.4), Inches(10.3), Inches(0.45),
            "DEVGTM MESH × LUMEN ANALYTICS", size=16, bold=True,
            color=TEAL, align=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(1.9), Inches(10.3), Inches(1.4),
            "Why now", size=64, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

add_textbox(sl, Inches(2), Inches(3.4), Inches(9.3), Inches(0.6),
            "A pitch tailored to your team, your stack, your moment",
            size=20, color=MUTED, align=PP_ALIGN.CENTER)

# Meta info row
meta = [
    ("Prospect", "Lumen Analytics"),
    ("Stage", "Series B · $42M"),
    ("Buyer", "VP RevOps"),
    ("Team", "12 SDRs · 18 AEs"),
    ("Stack", "HubSpot · Salesforce · Outreach"),
]
for i, (label, val) in enumerate(meta):
    x = Inches(0.3 + i * 2.55)
    add_rect(sl, x, Inches(4.2), Inches(2.45), Inches(1.4), CARD, BORDER)
    add_textbox(sl, x + Inches(0.15), Inches(4.3), Inches(2.15), Inches(0.35),
                label.upper(), size=9, bold=True, color=MUTED)
    add_textbox(sl, x + Inches(0.15), Inches(4.65), Inches(2.15), Inches(0.7),
                val, size=13, bold=True, color=TEXT)

slide_number(sl, 1)

# ──────────────────────────────────────────────────────────
# SLIDE 02 — Why we're reaching out
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Opening Hook", PURPLE_L)

add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.65),
            "We noticed something specific.", size=36, bold=True, color=TEXT)
add_textbox(sl, Inches(0.4), Inches(1.82), Inches(12.5), Inches(0.5),
            "We are not reaching out with a template. Three days ago, we came across something "
            "your CTO posted publicly on LinkedIn.",
            size=17, color=MUTED)

# Quote card
add_rect(sl, Inches(0.4), Inches(2.5), Inches(12.5), Inches(2.8), CARD, PURPLE)
add_textbox(sl, Inches(0.6), Inches(2.62), Inches(3.5), Inches(0.4),
            "PUBLIC POST · LINKEDIN · 9 DAYS AGO", size=11, bold=True, color=PURPLE_L)
add_textbox(sl, Inches(0.6), Inches(3.1), Inches(12), Inches(1.4),
            '"Manual prospect research is our SDR team\'s #1 bottleneck — we\'re losing hours a day."',
            size=24, bold=False, italic=True, color=TEXT)
add_rect(sl, Inches(0.4), Inches(5.2), Inches(0.06), Inches(0.7), PURPLE_L)
add_textbox(sl, Inches(0.6), Inches(5.25), Inches(4), Inches(0.4),
            "CTO, Lumen Analytics", size=14, bold=True, color=TEXT)
add_textbox(sl, Inches(0.6), Inches(5.65), Inches(6), Inches(0.35),
            "San Francisco · 180 employees · Data Observability", size=12, color=MUTED)

add_textbox(sl, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.5),
            "We build software that eliminates exactly this bottleneck — and we can show you the math on what it costs Lumen right now.",
            size=15, color=MUTED)
slide_number(sl, 2)

# ──────────────────────────────────────────────────────────
# SLIDE 03 — Signal grid (6 tiles)
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Signal Intelligence", TEAL)

add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.55),
            "What our AI saw across 6 sources", size=30, bold=True, color=TEXT)

signals = [
    ("Reo.Dev", "–14 days", "3 anonymous visits to getreo.dev/pricing",
     "Unique sessions on pricing page — active evaluation, not casual research.",
     "Buying signal · High confidence", PURPLE_L),
    ("Lumen Careers", "–9 days", '"Outbound SDR Manager" role opened',
     "Headcount investment in outbound — scaling the bottleneck the CTO flagged.",
     "Scaling signal · High confidence", AMBER),
    ("LinkedIn · Public", "–9 days", "CTO post: #1 bottleneck — losing hours/day",
     "Self-described pain, executive level, public record. Highest-signal statement possible.",
     "Pain signal · Very high confidence", BLUE_L),
    ("Crunchbase", "–6 months", "Series B closed · $42M raised",
     "Post-Series B companies have budget and board pressure to show GTM efficiency.",
     "Budget signal · High confidence", TEAL),
    ("Fireflies", "–30 days", "2 meetings with RevOps consultant (competitive)",
     "Lumen is actively evaluating alternatives — buying window is open now.",
     "Competitive signal · Medium confidence", PINK),
    ("Firecrawl", "–recent", 'Homepage hero: "AI-native data observability"',
     "Lumen updated their top-of-fold messaging toward AI positioning.",
     "Positioning signal · Medium confidence", GREEN),
]

cols = 3
for i, (source, ts, body, detail, conf, col) in enumerate(signals):
    row, c = divmod(i, cols)
    x = Inches(0.3 + c * 4.35)
    y = Inches(1.85 + row * 2.4)
    add_rect(sl, x, y, Inches(4.15), Inches(2.2), CARD, col)
    add_textbox(sl, x + Inches(0.15), y + Inches(0.12), Inches(2.5), Inches(0.35),
                source.upper(), size=10, bold=True, color=col)
    add_textbox(sl, x + Inches(3), y + Inches(0.12), Inches(1.0), Inches(0.35),
                ts, size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    add_textbox(sl, x + Inches(0.15), y + Inches(0.5), Inches(3.85), Inches(0.55),
                body, size=14, bold=True, color=TEXT)
    add_textbox(sl, x + Inches(0.15), y + Inches(1.1), Inches(3.85), Inches(0.55),
                detail, size=12, color=BODY)
    add_textbox(sl, x + Inches(0.15), y + Inches(1.72), Inches(3.85), Inches(0.35),
                "✦  " + conf, size=11, bold=True, color=col)
slide_number(sl, 3)

# ──────────────────────────────────────────────────────────
# SLIDE 04 — The pattern these signals form
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "AI Analysis", PURPLE_L)

add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.55),
            "Connect the dots", size=34, bold=True, color=TEXT)
add_textbox(sl, Inches(0.4), Inches(1.72), Inches(12.5), Inches(0.45),
            "Six signals. One coherent picture: Lumen is scaling outbound, "
            "acknowledging the bottleneck, holding budget, and actively evaluating solutions.",
            size=16, color=MUTED)

patterns = [
    (AMBER, "01", "Scaling Into the Problem",
     "You hired an Outbound SDR Manager 9 days ago — growing the team "
     "that is already spending 2.5 hours a day on manual research.",
     "source: Lumen Careers · job posting · April 2025"),
    (BLUE_L, "02", "Executive Acknowledgment",
     "The CTO publicly named manual research as the #1 SDR bottleneck. "
     "That is a self-report from the top of the org, on the record.",
     "source: LinkedIn public post · CTO · April 2025"),
    (TEAL, "03", "Budget + Active Evaluation",
     "Series B closed 6 months ago (budget exists). Two meetings with a RevOps "
     "consultant (evaluating). Three pricing-page visits (looked at us specifically).",
     "source: Crunchbase · Fireflies meeting logs · Reo.Dev"),
]

for pi, (col, num, title, body, signal) in enumerate(patterns):
    y = Inches(2.38 + pi * 1.38)
    add_rect(sl, Inches(0.4), y, Inches(12.5), Inches(1.25), CARD, col)
    add_rect(sl, Inches(0.4), y, Inches(0.9), Inches(1.25), CARD, col)
    add_textbox(sl, Inches(0.4), y + Inches(0.2), Inches(0.9), Inches(0.8),
                num, size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(sl, Inches(1.45), y + Inches(0.08), Inches(4.5), Inches(0.4),
                title.upper(), size=11, bold=True, color=col)
    add_textbox(sl, Inches(1.45), y + Inches(0.45), Inches(8.0), Inches(0.55),
                body, size=14, color=TEXT)
    add_textbox(sl, Inches(1.45), y + Inches(0.97), Inches(10.8), Inches(0.25),
                signal, size=11, color=MUTED)

slide_number(sl, 4)

# ──────────────────────────────────────────────────────────
# SLIDE 05 — CTO Quote (visual centerpiece)
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Your Stated Bottleneck", PURPLE_L)

# Large quote mark
add_textbox(sl, Inches(1.5), Inches(0.9), Inches(2), Inches(1.5),
            "“", size=120, bold=True, color=PURPLE_L)

# The quote — large and centered
add_textbox(sl, Inches(1.2), Inches(2.0), Inches(11), Inches(2.2),
            "Manual prospect research is our SDR team’s "
            "#1 bottleneck — we’re losing hours a day.",
            size=38, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

# Attribution block
add_rect(sl, Inches(4.5), Inches(4.5), Inches(4.3), Inches(0.06), PURPLE_L)
add_textbox(sl, Inches(4.0), Inches(4.7), Inches(5.3), Inches(0.45),
            "CTO, Lumen Analytics", size=20, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(3.5), Inches(5.18), Inches(6.3), Inches(0.4),
            "Data Observability · Series B · San Francisco", size=16, color=MUTED, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(3.8), Inches(5.6), Inches(5.7), Inches(0.35),
            "LinkedIn public post · April 17, 2025", size=14, color=PURPLE_L, align=PP_ALIGN.CENTER)

# AI attribution note
add_rect(sl, Inches(3.2), Inches(6.3), Inches(6.9), Inches(0.75), CARD, PURPLE)
add_textbox(sl, Inches(3.2), Inches(6.38), Inches(6.9), Inches(0.58),
            "The AI Pitch Generator surfaced this verbatim quote from LinkedIn's "
            "public feed and flagged it as the highest-priority outreach signal in Lumen's account profile.",
            size=12, color=MUTED, align=PP_ALIGN.CENTER)
slide_number(sl, 5)

# ──────────────────────────────────────────────────────────
# SLIDE 06 — Cost math (formula)
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Cost Analysis", AMBER)

add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.6),
            "What this bottleneck costs Lumen — annually", size=30, bold=True, color=TEXT)

# Formula box
add_rect(sl, Inches(0.4), Inches(1.9), Inches(12.5), Inches(4.0), CARD, PURPLE)

formula_rows = [
    ("SDR headcount", "=", "12 SDRs", "from firmographic enrichment"),
    ("Daily research time per SDR", "×", "2.5 hrs / day", "industry average · data obs. SaaS"),
    ("Fully-loaded SDR cost", "×", "$65 / hr", "salary + benefits + tools"),
    ("Working days per year", "×", "220 days", "excluding PTO and holidays"),
]
for ri, (label, op, val, note) in enumerate(formula_rows):
    y = Inches(2.05 + ri * 0.72)
    add_textbox(sl, Inches(0.65), y, Inches(3.6), Inches(0.6), label, size=14, color=BODY)
    add_textbox(sl, Inches(4.3), y, Inches(0.5), Inches(0.6), op, size=22, bold=True, color=PURPLE_L, align=PP_ALIGN.CENTER)
    add_textbox(sl, Inches(4.95), y, Inches(2.2), Inches(0.6), val, size=18, bold=True, color=TEXT)
    add_textbox(sl, Inches(7.3), y, Inches(5.4), Inches(0.6), note, size=12, color=MUTED)

# Divider
add_rect(sl, Inches(0.55), Inches(4.97), Inches(12.2), Inches(0.02), BORDER)

# Result row
add_textbox(sl, Inches(0.65), Inches(5.1), Inches(4.5), Inches(0.65),
            "Annual research cost:", size=18, bold=True, color=MUTED)
add_textbox(sl, Inches(5.1), Inches(4.95), Inches(4.5), Inches(0.9),
            "~$429,000", size=48, bold=True, color=TEAL)
add_textbox(sl, Inches(9.7), Inches(5.3), Inches(3), Inches(0.5),
            "12 × 2.5 × $65 × 220 = $429,000", size=13, color=MUTED)

add_textbox(sl, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.4),
            "These are Lumen’s numbers. Our AI applied the 12-SDR headcount "
            "to the industry time benchmark.",
            size=13, color=MUTED)
slide_number(sl, 6)

# ──────────────────────────────────────────────────────────
# SLIDE 07 — Industry benchmark
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Industry Benchmark", AMBER)

add_textbox(sl, Inches(4.0), Inches(1.1), Inches(5.3), Inches(1.6),
            "38%", size=96, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(2.0), Inches(2.75), Inches(9.3), Inches(0.6),
            "of SDR time spent on manual prospect research",
            size=22, color=TEXT, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(2.0), Inches(3.38), Inches(9.3), Inches(0.4),
            "data observability SaaS companies at Lumen’s stage (illustrative benchmark)",
            size=15, color=MUTED, align=PP_ALIGN.CENTER)

bench_cards = [
    ("🔍", "38%", AMBER, "Lumen estimated research share\n(industry average applied)"),
    ("⚡", "62%", PURPLE_L, "Remaining SDR capacity\nfor selling activities"),
    ("🎯", ">90%", TEAL, "Target selling-time ratio\nwith DevGTM automation"),
]
for bi, (icon, val, col, desc) in enumerate(bench_cards):
    x = Inches(0.6 + bi * 4.1)
    add_rect(sl, x, Inches(4.1), Inches(3.9), Inches(2.5), CARD, col)
    add_textbox(sl, x + Inches(0.2), Inches(4.22), Inches(3.5), Inches(0.5),
                icon, size=28)
    add_textbox(sl, x + Inches(0.2), Inches(4.75), Inches(3.5), Inches(0.85),
                val, size=44, bold=True, color=col)
    add_textbox(sl, x + Inches(0.2), Inches(5.6), Inches(3.5), Inches(0.85),
                desc, size=14, color=BODY)
slide_number(sl, 7)

# ──────────────────────────────────────────────────────────
# SLIDE 08 — What DevGTM Mesh does (3 pillars)
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "DevGTM Mesh", PURPLE_L)

add_textbox(sl, Inches(1.5), Inches(1.05), Inches(10.3), Inches(0.75),
            "AI-powered GTM intelligence built for developer-led sales",
            size=28, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

pillars = [
    (PURPLE_L, "01", "🔗", "Unify your GTM data",
     "HubSpot, Salesforce, Outreach, Reo.Dev, Fireflies, Firecrawl — all normalized into one "
     "account record. No silos. No manual reconciliation."),
    (TEAL, "02", "⚡", "Score intent automatically",
     "Six signal sources weighted and aggregated into a single intent score per account. "
     "Your reps know which accounts to call today — before they open Salesforce."),
    (AMBER, "03", "✦", "Generate personalized pitches",
     "The same way this deck was generated: AI reads your prospect's signals, writes the pitch, "
     "and delivers it as a shareable URL. Not a template — a reasoned argument."),
]
for pi, (col, num, icon, title, desc) in enumerate(pillars):
    x = Inches(0.4 + pi * 4.3)
    add_rect(sl, x, Inches(2.0), Inches(4.1), Inches(5.0), CARD, col)
    add_textbox(sl, x + Inches(0.2), Inches(2.1), Inches(3.7), Inches(0.6),
                num, size=44, bold=True, color=col)
    add_textbox(sl, x + Inches(0.2), Inches(2.75), Inches(3.7), Inches(0.5),
                icon + "  " + title, size=20, bold=True, color=col)
    add_textbox(sl, x + Inches(0.2), Inches(3.35), Inches(3.7), Inches(3.5),
                desc, size=15, color=BODY)
slide_number(sl, 8)

# ──────────────────────────────────────────────────────────
# SLIDE 09 — Stack fit
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Stack Fit · Lumen’s Exact Tools", TEAL)

add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.6),
            "Native integrations with your current stack", size=32, bold=True, color=TEXT)

tools = [
    ("🟠", "HubSpot", "CRM", RGBColor(0xff, 0x7a, 0x59),
     "Native connector syncs companies, contacts, deals, and activities in real time."),
    ("☁️", "Salesforce", "CRM · Source of truth", BLUE_L,
     "Bidirectional sync: intent scores and research briefs visible in SF records."),
    ("🟣", "Outreach", "Sales Engagement", PURPLE_L,
     "AI-generated pitch content flows directly into Outreach sequences."),
    ("⬡", "DevGTM Mesh", "Intelligence Layer", TEAL,
     "Sits above your stack. Reads from all three. Writes enrichment and scores back."),
]
for ti, (icon, name, role, col, desc) in enumerate(tools):
    x = Inches(0.3 + ti * 3.25)
    if ti < len(tools) - 1:
        add_rect(sl, x, Inches(1.9), Inches(3.05), Inches(3.5), CARD, col)
    else:
        # DevGTM card — featured
        add_rect(sl, x, Inches(1.9), Inches(3.05), Inches(3.5), CARD, PURPLE)
    add_textbox(sl, x + Inches(0.15), Inches(2.0), Inches(2.75), Inches(0.5),
                icon + "  " + name, size=20, bold=True, color=col)
    add_textbox(sl, x + Inches(0.15), Inches(2.55), Inches(2.75), Inches(0.35),
                role.upper(), size=10, bold=True, color=MUTED)
    add_textbox(sl, x + Inches(0.15), Inches(2.98), Inches(2.75), Inches(2.25),
                desc, size=14, color=BODY)
    if ti < len(tools) - 1:
        add_textbox(sl, x + Inches(3.05), Inches(2.95), Inches(0.2), Inches(0.55),
                    "+", size=24, bold=True, color=BORDER, align=PP_ALIGN.CENTER)

# Callout bar
add_rect(sl, Inches(0.3), Inches(5.6), Inches(12.73), Inches(0.8), CARD, TEAL)
add_textbox(sl, Inches(0.45), Inches(5.7), Inches(12.43), Inches(0.6),
            "No rip-and-replace. DevGTM Mesh enhances your existing HubSpot + Salesforce + Outreach "
            "investment — it does not replace any of them.",
            size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
slide_number(sl, 9)

# ──────────────────────────────────────────────────────────
# HELPER: Before / After table slide
# ──────────────────────────────────────────────────────────
def ba_slide(slide_n, section_lbl, section_col, h_title, rows_before, rows_after,
             before_footer, after_footer):
    sl = add_slide(); fill_bg(sl)
    add_logo(sl)
    section_label(sl, section_lbl, section_col)
    add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.6),
                h_title, size=28, bold=True, color=TEXT)

    # Before column
    add_rect(sl, Inches(0.4), Inches(1.9), Inches(6.0), Inches(0.55), CARD, RED_L)
    add_textbox(sl, Inches(0.55), Inches(1.96), Inches(5.7), Inches(0.43),
                "✗  Today — Manual", size=16, bold=True, color=RED_L)
    for ri, (icon, text, timing) in enumerate(rows_before):
        y = Inches(2.55 + ri * 0.82)
        add_rect(sl, Inches(0.4), y, Inches(6.0), Inches(0.75), CARD, BORDER)
        add_textbox(sl, Inches(0.55), y + Inches(0.04), Inches(0.4), Inches(0.65), icon, size=18)
        add_textbox(sl, Inches(1.05), y + Inches(0.04), Inches(5.15), Inches(0.42), text, size=14, color=BODY)
        if timing:
            add_textbox(sl, Inches(1.05), y + Inches(0.45), Inches(3.5), Inches(0.28),
                        timing, size=11, color=MUTED)
    add_rect(sl, Inches(0.4), Inches(6.1), Inches(6.0), Inches(0.65), CARD, RED_L)
    add_textbox(sl, Inches(0.55), Inches(6.18), Inches(5.7), Inches(0.5),
                before_footer, size=13, bold=True, color=RED_L)

    # After column
    add_rect(sl, Inches(6.7), Inches(1.9), Inches(6.2), Inches(0.55), CARD, GREEN_L)
    add_textbox(sl, Inches(6.85), Inches(1.96), Inches(5.9), Inches(0.43),
                "✓  With DevGTM Mesh", size=16, bold=True, color=GREEN_L)
    for ri, (icon, text, timing) in enumerate(rows_after):
        y = Inches(2.55 + ri * 0.82)
        add_rect(sl, Inches(6.7), y, Inches(6.2), Inches(0.75), CARD, BORDER)
        add_textbox(sl, Inches(6.85), y + Inches(0.04), Inches(0.4), Inches(0.65), icon, size=18)
        add_textbox(sl, Inches(7.35), y + Inches(0.04), Inches(5.35), Inches(0.42), text, size=14, color=BODY)
        if timing:
            add_textbox(sl, Inches(7.35), y + Inches(0.45), Inches(3.5), Inches(0.28),
                        timing, size=11, color=MUTED)
    add_rect(sl, Inches(6.7), Inches(6.1), Inches(6.2), Inches(0.65), CARD, GREEN_L)
    add_textbox(sl, Inches(6.85), Inches(6.18), Inches(5.9), Inches(0.5),
                after_footer, size=13, bold=True, color=GREEN_L)

    slide_number(sl, slide_n)

# ──────────────────────────────────────────────────────────
# SLIDE 10 — SDR Before / After
# ──────────────────────────────────────────────────────────
ba_slide(
    10,
    "12 SDRs · Daily Workflow", AMBER,
    "SDR research workflow — before & after",
    [
        ("🔍", "Open Google, LinkedIn, company site — 5+ tabs per account", "~25 min / account"),
        ("📋", "Manually copy company facts into CRM / notepad", "~10 min / account"),
        ("✍️", 'Write a "personalized" opening from generic template', "~15 min / account"),
        ("❌", "No signal awareness — outreach is time-based, not intent-based", None),
    ],
    [
        ("⚡", "Intent-ranked account queue surfaces in dashboard — no prospecting list", "~0 min"),
        ("📊", "Research brief pre-generated: company summary, tech stack, contacts", "~2 min review"),
        ("✦", "AI-generated pitch from account signals — contextually justified", "~3 min review + send"),
        ("🎯", "Signal-triggered outreach: SDR knows why this account, why today", None),
    ],
    "Total: ~2.5 hrs/day on research = $429K/yr across 12 SDRs",
    "Total: ~5 min / account · Research time recoverable: ~$200K/yr",
)

# ──────────────────────────────────────────────────────────
# SLIDE 11 — AE Before / After
# ──────────────────────────────────────────────────────────
ba_slide(
    11,
    "18 AEs · Meeting Prep Workflow", TEAL,
    "AE meeting prep — before & after",
    [
        ("📂", "Re-read 3+ CRM notes, old emails, Fireflies transcripts separately", "~30 min / meeting"),
        ("🔎", "Look up recent news, job postings, company announcements manually", "~15 min / meeting"),
        ("🗒️", "Write call prep notes with no signal context", "~15 min / meeting"),
        ("❌", "No awareness of competitive activity or intent-score changes since last call", None),
    ],
    [
        ("📋", "Pre-meeting brief auto-assembled: prior summaries, action items, sentiment", "~2 min review"),
        ("📈", 'Intent score delta shown: "account spiked 23 pts this week — here is why"', "~1 min"),
        ("🎙️", "Fireflies transcript summary + action items surfaced inline — one click", "~2 min"),
        ("🏆", '"Competitor meeting logged 11 days ago — come prepared" flagged automatically', None),
    ],
    "Total: ~1 hr/meeting on manual prep · 3–5 meetings/week/AE",
    "Total: ~5 min / meeting prep · context is pre-loaded before the call",
)

# ──────────────────────────────────────────────────────────
# SLIDE 12 — Case study: Sentinel Metrics
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Customer Proof", PURPLE_L)

add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.55),
            "Sentinel Metrics — analogous customer profile", size=30, bold=True, color=TEXT)

# Left card
add_rect(sl, Inches(0.4), Inches(1.85), Inches(6.2), Inches(5.4), CARD, PURPLE)
add_textbox(sl, Inches(0.6), Inches(1.98), Inches(5.8), Inches(0.7),
            "Sentinel Metrics", size=32, bold=True, color=TEXT)
add_textbox(sl, Inches(0.6), Inches(2.72), Inches(5.8), Inches(0.55),
            "Data observability SaaS · 200 employees · Series B · San Francisco\n"
            "HubSpot + Salesforce + Outreach · 14 SDRs · 20 AEs",
            size=14, color=MUTED)

metrics = [
    ("47%", TEAL, "reduction in SDR research time in first 90 days"),
    ("2.3×", PURPLE_L, "pipeline increase in 90 days — same team size"),
    ("90", AMBER, "days from signed contract to measurable pipeline impact"),
]
for mi, (val, col, desc) in enumerate(metrics):
    y = Inches(3.4 + mi * 1.1)
    add_textbox(sl, Inches(0.6), y, Inches(1.8), Inches(0.9), val, size=40, bold=True, color=col)
    add_textbox(sl, Inches(2.5), y + Inches(0.15), Inches(3.8), Inches(0.7), desc, size=15, color=BODY)

# Right side
add_rect(sl, Inches(6.85), Inches(1.85), Inches(6.1), Inches(3.0), CARD, PURPLE)
add_textbox(sl, Inches(7.05), Inches(1.97), Inches(5.7), Inches(2.7),
            "“We went from SDRs spending half their day on research to having a "
            "morning queue pre-loaded with context. The first month was setup. "
            "By month three, pipeline was unrecognizable.”",
            size=16, italic=True, color=TEXT)
add_textbox(sl, Inches(7.05), Inches(4.7), Inches(5.7), Inches(0.4),
            "— VP Sales, Sentinel Metrics · post-deployment review",
            size=12, bold=True, color=PURPLE_L)

add_rect(sl, Inches(6.85), Inches(5.1), Inches(6.1), Inches(2.15), CARD, BORDER)
add_textbox(sl, Inches(7.05), Inches(5.2), Inches(5.7), Inches(0.4),
            "WHY THIS COMPARISON IS RELEVANT", size=11, bold=True, color=MUTED)
add_textbox(sl, Inches(7.05), Inches(5.65), Inches(5.7), Inches(1.45),
            "Sentinel’s profile at deployment closely matched Lumen today: same industry, "
            "same team scale, same GTM stack, same CTO-acknowledged research bottleneck.",
            size=14, color=BODY)
slide_number(sl, 12)

# ──────────────────────────────────────────────────────────
# SLIDE 13 — Projected impact for Lumen
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Lumen-Specific Projection", TEAL)

add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.6),
            "Applying Sentinel ratios to Lumen’s 12-SDR / 18-AE team",
            size=28, bold=True, color=TEXT)

proj_cards = [
    ("⏱️", "~$200K", TEAL, "Annual research time recovered", "47% × $429K baseline (slide 6)"),
    ("+📈", "+$3.2M", PURPLE_L, "Estimated pipeline uplift at current ACV", "2.3× ratio · $185K ACV · current conversion"),
    ("🎯", "90", AMBER, "Days to measurable pipeline impact", "based on Sentinel deployment timeline"),
    ("👥", "30", GREEN, "People impacted (12 SDRs + 18 AEs)", "every rep gets a pre-loaded research queue"),
]
for pi, (icon, val, col, label, basis) in enumerate(proj_cards):
    x = Inches(0.4 + pi * 3.25)
    add_rect(sl, x, Inches(2.0), Inches(3.05), Inches(3.5), CARD, col)
    add_textbox(sl, x + Inches(0.2), Inches(2.1), Inches(2.65), Inches(0.55), icon, size=28)
    add_textbox(sl, x + Inches(0.2), Inches(2.7), Inches(2.65), Inches(0.85),
                val, size=40, bold=True, color=col)
    add_textbox(sl, x + Inches(0.2), Inches(3.6), Inches(2.65), Inches(0.75),
                label, size=14, color=BODY)
    add_textbox(sl, x + Inches(0.2), Inches(4.35), Inches(2.65), Inches(0.65),
                basis, size=11, color=MUTED)

# Note
add_rect(sl, Inches(0.4), Inches(5.75), Inches(12.5), Inches(0.9), CARD, TEAL)
add_textbox(sl, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.7),
            "These are AI-generated projections based on Sentinel Metrics outcomes applied to Lumen’s known "
            "team structure. They are directional estimates — actual results depend on your pipeline volume and conversion rates.",
            size=13, color=MUTED, align=PP_ALIGN.CENTER)
slide_number(sl, 13)

# ──────────────────────────────────────────────────────────
# SLIDE 14 — 30/60/90 rollout plan
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Implementation Plan · Lumen Stack", PURPLE_L)

add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.6),
            "30 / 60 / 90 day rollout — your stack, your milestones",
            size=28, bold=True, color=TEXT)

phases = [
    (AMBER, "Phase 1", "Days 1–30", [
        ("Wk 1", "HubSpot connector activated — companies, contacts, deals syncing in real time"),
        ("Wk 1", "Salesforce bidirectional sync configured — intent scores visible in SF records"),
        ("Wk 2", "Outreach integration live — AI pitch content flows into sequences"),
        ("Wk 3–4", "SDR onboarding (2 sessions) · first AI-generated research briefs reviewed"),
        ("Day 30", "Milestone check: research time per account, SDR feedback, pipeline delta"),
    ]),
    (PURPLE_L, "Phase 2", "Days 31–60", [
        ("Wk 5", "Intent score calibration — adjust signal weights to match Lumen’s ICP patterns"),
        ("Wk 6", "AE meeting prep workflows enabled — Fireflies summaries surfaced in-platform"),
        ("Wk 7", "Firecrawl enrichment pipeline running on full account list — homepage signals active"),
        ("Wk 8", "Pitch Generator enabled for all SDRs · shareable pitch URL workflow trained"),
        ("Day 60", "Milestone check: pipeline created, pitch send-rate vs. baseline, research time delta"),
    ]),
    (TEAL, "Phase 3", "Days 61–90", [
        ("Wk 9", "Full GTM chat interface live — VP RevOps using natural-language queries"),
        ("Wk 10", "SDR auto-assignment rules configured — accounts routed by territory + intent score"),
        ("Wk 11", "Reo.Dev developer-intent signals layered in — pricing-page visit alerts live"),
        ("Wk 12", "Full-team adoption review · success metrics report prepared for board"),
        ("Day 90", "Pilot success review: research time, pipeline, ACV impact vs. targets"),
    ]),
]
for ci, (col, phase, period, milestones) in enumerate(phases):
    x = Inches(0.3 + ci * 4.37)
    # Phase header
    add_rect(sl, x, Inches(1.9), Inches(4.17), Inches(0.65), CARD, col)
    add_textbox(sl, x + Inches(0.12), Inches(1.95), Inches(1.8), Inches(0.33),
                phase.upper(), size=11, bold=True, color=col)
    add_textbox(sl, x + Inches(0.12), Inches(2.3), Inches(1.8), Inches(0.3),
                period, size=13, bold=True, color=TEXT)
    # Milestones
    for mi, (week, task) in enumerate(milestones):
        y = Inches(2.65 + mi * 0.92)
        add_rect(sl, x, y, Inches(4.17), Inches(0.85), CARD, BORDER)
        add_textbox(sl, x + Inches(0.1), y + Inches(0.05), Inches(0.7), Inches(0.75),
                    "✦", size=14, color=col)
        add_textbox(sl, x + Inches(0.35), y + Inches(0.05), Inches(0.8), Inches(0.3),
                    week, size=11, bold=True, color=MUTED)
        add_textbox(sl, x + Inches(0.35), y + Inches(0.32), Inches(3.7), Inches(0.5),
                    task, size=12, color=BODY)

slide_number(sl, 14)

# ──────────────────────────────────────────────────────────
# SLIDE 15 — Investment (pricing)
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Investment", PURPLE_L)

add_textbox(sl, Inches(1.5), Inches(0.9), Inches(10.3), Inches(0.6),
            "Simple, transparent pricing", size=32, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(2), Inches(1.52), Inches(9.3), Inches(0.4),
            "Two options — both designed to return investment within 90 days",
            size=16, color=MUTED, align=PP_ALIGN.CENTER)

pricing_cards = [
    (False, PURPLE_L, "Seat-Based", "Per Rep", "Pricing on request · per seat / month", [
        "Scale with your team — pay for active seats only",
        "All integrations included (HubSpot + Salesforce + Outreach)",
        "AI Research Bot, Pitch Generator, Intent Scoring",
        "Best for growing teams expecting headcount changes",
    ]),
    (True, TEAL, "Recommended for Lumen", "Platform", "Pricing on request · flat platform rate / month", [
        "Predictable cost regardless of headcount growth",
        "All features + unlimited seats for your 30-person GTM team",
        "Dedicated onboarding and calibration support",
        "Recommended: Lumen is actively hiring (SDR Manager role open)",
    ]),
]
for pi, (featured, col, badge, name, price, features) in enumerate(pricing_cards):
    x = Inches(0.6 + pi * 6.3)
    border_col = col if featured else BORDER
    add_rect(sl, x, Inches(2.1), Inches(6.0), Inches(4.5), CARD, border_col)
    add_textbox(sl, x + Inches(0.25), Inches(2.2), Inches(5.5), Inches(0.4),
                badge.upper(), size=11, bold=True, color=col)
    add_textbox(sl, x + Inches(0.25), Inches(2.62), Inches(5.5), Inches(0.65),
                name, size=30, bold=True, color=col)
    add_textbox(sl, x + Inches(0.25), Inches(3.3), Inches(5.5), Inches(0.35),
                price, size=14, color=MUTED)
    for fi, feat in enumerate(features):
        add_textbox(sl, x + Inches(0.25), Inches(3.78 + fi * 0.6), Inches(5.5), Inches(0.55),
                    "✔  " + feat, size=14, color=BODY)

# ROI bar
add_rect(sl, Inches(0.6), Inches(6.75), Inches(12.13), Inches(0.55), CARD, TEAL)
add_textbox(sl, Inches(0.75), Inches(6.83), Inches(11.83), Inches(0.43),
            "ROI in under 90 days — the $200K/yr time recovery alone covers the platform cost within the first quarter",
            size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
slide_number(sl, 15)

# ──────────────────────────────────────────────────────────
# SLIDE 16 — Risk reversal
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Risk Reversal", TEAL)

# Hero box
add_rect(sl, Inches(0.5), Inches(1.0), Inches(12.33), Inches(5.2),
         CARD, PURPLE)
add_textbox(sl, Inches(1), Inches(1.2), Inches(11.33), Inches(1.2),
            "60-Day Paid Pilot with Full Refund Guarantee",
            size=36, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1.5), Inches(2.45), Inches(10.33), Inches(0.75),
            "We define success metrics together before day one. If we do not hit them by day 60, "
            "you get your money back — no conditions, no questions.",
            size=18, color=MUTED, align=PP_ALIGN.CENTER)

risk_terms = [
    ("📅", "60", "Days in the pilot window"),
    ("📋", "Day 1", "Success metrics defined before any work begins"),
    ("💰", "100%", "Refund if targets are missed"),
    ("⚡", "Wk 1", "HubSpot + Salesforce + Outreach integration live"),
]
for ri, (icon, val, desc) in enumerate(risk_terms):
    x = Inches(0.8 + ri * 3.1)
    add_rect(sl, x, Inches(3.35), Inches(2.8), Inches(2.65), CARD, TEAL)
    add_textbox(sl, x + Inches(0.2), Inches(3.45), Inches(2.4), Inches(0.5), icon, size=26)
    add_textbox(sl, x + Inches(0.2), Inches(3.98), Inches(2.4), Inches(0.7),
                val, size=28, bold=True, color=TEAL)
    add_textbox(sl, x + Inches(0.2), Inches(4.7), Inches(2.4), Inches(0.85),
                desc, size=13, color=BODY, align=PP_ALIGN.CENTER)

# Guarantee text
add_textbox(sl, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.7),
            "Suggested success metrics: (1) >30% reduction in SDR research time, "
            "(2) >$500K pipeline in 60 days attributable to AI-prioritized outreach, "
            "(3) >80% SDR adoption rate. Negotiable — we define them together.",
            size=13, color=MUTED, align=PP_ALIGN.CENTER)
slide_number(sl, 16)

# ──────────────────────────────────────────────────────────
# SLIDE 17 — Next step
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Recommended Next Step", PURPLE_L)

add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.6),
            "One 30-minute working session — not a demo",
            size=30, bold=True, color=TEXT)

# Ask card (left)
add_rect(sl, Inches(0.4), Inches(1.9), Inches(6.2), Inches(5.3), CARD, PURPLE)
add_textbox(sl, Inches(0.6), Inches(2.0), Inches(5.8), Inches(0.4),
            "THE ASK", size=11, bold=True, color=PURPLE_L)
add_textbox(sl, Inches(0.6), Inches(2.42), Inches(5.8), Inches(0.8),
            "30-minute working session", size=26, bold=True, color=TEXT)
add_textbox(sl, Inches(0.6), Inches(3.26), Inches(5.8), Inches(0.35),
            "Suggested attendees from your team:", size=14, color=MUTED)
attendees = [
    "VP RevOps — to validate the cost math and success metric framing",
    "1 SDR Manager — to walk through the daily workflow change (slide 10)",
    "1 Account Executive — to review the meeting prep workflow (slide 11)",
]
for ai, att in enumerate(attendees):
    add_textbox(sl, Inches(0.75), Inches(3.7 + ai * 0.55), Inches(5.55), Inches(0.48),
                "●  " + att, size=13, color=BODY)

# Date options
add_textbox(sl, Inches(0.6), Inches(5.55), Inches(5.8), Inches(0.35),
            "Proposed dates (flexible):", size=13, color=MUTED)
for di, (label, val) in enumerate([("Option A", "May 6, 2025 · 10:00 AM PT"), ("Option B", "May 8, 2025 · 2:00 PM PT")]):
    x = Inches(0.6 + di * 2.9)
    add_rect(sl, x, Inches(5.98), Inches(2.7), Inches(0.9), CARD, PURPLE)
    add_textbox(sl, x + Inches(0.1), Inches(6.05), Inches(2.5), Inches(0.28),
                label.upper(), size=10, bold=True, color=PURPLE_L)
    add_textbox(sl, x + Inches(0.1), Inches(6.35), Inches(2.5), Inches(0.45),
                val, size=12, color=TEXT)

# Agenda card (right)
add_rect(sl, Inches(6.85), Inches(1.9), Inches(6.1), Inches(5.3), CARD, BORDER)
add_textbox(sl, Inches(7.05), Inches(2.0), Inches(5.7), Inches(0.4),
            "30-MINUTE AGENDA", size=11, bold=True, color=MUTED)
agenda = [
    ("01", "Validate the cost math (slide 6) against Lumen’s actual SDR research time — "
           "do the numbers hold? (10 min)"),
    ("02", "Walk the SDR manager through the before/after workflow live in the platform — "
           "hands on, no slides (12 min)"),
    ("03", "Define the pilot success metrics together — agree on the numbers that would make "
           "this a clear yes after 60 days (8 min)"),
]
for ai, (num, text) in enumerate(agenda):
    y = Inches(2.55 + ai * 1.35)
    add_textbox(sl, Inches(7.05), y, Inches(0.55), Inches(0.35), num, size=14, bold=True, color=PURPLE_L)
    add_textbox(sl, Inches(7.65), y, Inches(5.1), Inches(1.2), text, size=14, color=BODY)

add_rect(sl, Inches(6.85), Inches(6.4), Inches(6.1), Inches(0.6), CARD, PURPLE)
add_textbox(sl, Inches(7.0), Inches(6.48), Inches(5.8), Inches(0.44),
            "Working session, not a sales call. Bring your actual numbers.",
            size=13, bold=True, color=PURPLE_L, align=PP_ALIGN.CENTER)
slide_number(sl, 17)

# ──────────────────────────────────────────────────────────
# SLIDE 18 — Appendix: data sources
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_label(sl, "Appendix · AI Audit Trail", PURPLE_L)

add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.55),
            "Every claim in this deck is traceable to a source", size=28, bold=True, color=TEXT)
add_textbox(sl, Inches(0.4), Inches(1.7), Inches(12.5), Inches(0.45),
            "The AI Pitch Generator assembled these signals and mapped each to the slides where it was used.",
            size=14, color=MUTED)

# Table header
col_widths = [Inches(0.75), Inches(3.8), Inches(3.2), Inches(2.8), Inches(1.7)]
headers = ["Slide", "Signal Used", "Source", "Date / Recency", "Confidence"]
add_rect(sl, Inches(0.3), Inches(2.28), Inches(12.73), Inches(0.5), CARD, BORDER)
x_cursor = Inches(0.3)
for hi, (hdr, cw) in enumerate(zip(headers, col_widths)):
    add_textbox(sl, x_cursor + Inches(0.1), Inches(2.33), cw - Inches(0.1), Inches(0.4),
                hdr.upper(), size=10, bold=True, color=MUTED)
    x_cursor += cw

# Table rows
rows = [
    ("2, 5",     "CTO public post: \"#1 bottleneck — losing hours a day\"",
     "LinkedIn · Public Feed Monitor",       "April 17, 2025 (–9 days)", ("Very High", GREEN)),
    ("3, 4",     "3 anonymous pricing-page visits · getreo.dev/pricing",
     "Reo.Dev · Visitor Intelligence",       "April 12–25, 2025 (–14 days)", ("High", GREEN)),
    ("3, 4, 15", "Job posting: \"Outbound SDR Manager\" opened",
     "Lumen Careers · Firecrawl",            "April 17, 2025 (–9 days)", ("High", GREEN)),
    ("3, 4, 6, 13", "Series B · $42M raised · 12 SDRs · 18 AEs · $185K ACV",
     "Crunchbase · Firecrawl enrichment",    "October 2024 (–6 months)", ("High", GREEN)),
    ("3, 4",     "2 meetings with RevOps consultant (competitive)",
     "Fireflies · Meeting Intelligence Sync","March–April 2025 (–30 days)", ("Medium", AMBER)),
    ("3",        "Homepage hero: \"AI-native data observability\"",
     "Firecrawl · lumenanalytics.com",       "Recent (detected on crawl)", ("Medium", AMBER)),
    ("9",        "GTM stack: HubSpot + Salesforce + Outreach confirmed",
     "HubSpot record · Job descriptions",    "Current (cold record)", ("High", GREEN)),
    ("6, 10, 13","SDR headcount (12), AE headcount (18), ACV ($185K)",
     "Firecrawl enrichment · LinkedIn data", "April 2025", ("High", GREEN)),
]
for ri, (slide_ref, signal, source, date, (conf_label, conf_col)) in enumerate(rows):
    y = Inches(2.82 + ri * 0.54)
    row_bg = CARD if ri % 2 == 0 else BG
    add_rect(sl, Inches(0.3), y, Inches(12.73), Inches(0.52), row_bg, BORDER, 0)
    x_cursor = Inches(0.3)
    cell_data = [
        (slide_ref, PURPLE_L, 10, True),
        (signal, TEXT, 12, False),
        (source, TEAL, 11, False),
        (date, MUTED, 11, False),
    ]
    for ci2, ((val, col, sz, bld), cw) in enumerate(zip(cell_data, col_widths)):
        add_textbox(sl, x_cursor + Inches(0.1), y + Inches(0.06), cw - Inches(0.15), Inches(0.42),
                    val, size=sz, bold=bld, color=col)
        x_cursor += cw
    add_textbox(sl, x_cursor + Inches(0.1), y + Inches(0.06), col_widths[4] - Inches(0.15), Inches(0.42),
                conf_label, size=12, bold=True, color=conf_col)
slide_number(sl, 18)

# ──────────────────────────────────────────────────────────
# Save
# ──────────────────────────────────────────────────────────
out_path = "/home/shekerk/cc-best-practice-clone/presentation/ai-pitch-generator-demo/ai-pitch-generator-demo.pptx"
prs.save(out_path)
print(f"Saved {len(prs.slides)} slides to {out_path}")
