"""
Generates devgtm-ai-use-cases.pptx — a 18-slide PowerPoint deck
matching the HTML presentation structure and dark-navy DevGTM color palette.
Run: python3 build-deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import sys

# ── Color palette ─────────────────────────────────────────
BG       = RGBColor(0x0b, 0x0f, 0x1a)   # dark navy
CARD     = RGBColor(0x13, 0x19, 0x29)   # card bg
BORDER   = RGBColor(0x1e, 0x2d, 0x47)   # borders
TEXT     = RGBColor(0xe8, 0xed, 0xf5)   # primary text
MUTED    = RGBColor(0x6b, 0x7f, 0xa3)   # muted text
BODY     = RGBColor(0xa8, 0xbd, 0xd4)   # body text

PURPLE   = RGBColor(0x7c, 0x3a, 0xed)   # primary purple
PURPLE_L = RGBColor(0xa7, 0x8b, 0xfa)   # light purple
TEAL     = RGBColor(0x2d, 0xd4, 0xbf)   # teal accent
AMBER    = RGBColor(0xf5, 0x9e, 0x0b)   # amber
BLUE     = RGBColor(0x3b, 0x82, 0xf6)   # blue
BLUE_L   = RGBColor(0x60, 0xa5, 0xfa)   # light blue
GREEN    = RGBColor(0x34, 0xd3, 0x99)   # green

# Category colors
CAT_INTEL = PURPLE_L
CAT_CONV  = TEAL
CAT_PRIO  = AMBER
CAT_DATA  = BLUE_L

# ── Slide dimensions (widescreen 16:9) ────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color=BG):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_logo(slide):
    """Small top-left logo mark."""
    add_rect(slide, Inches(0.4), Inches(0.25), Inches(0.35), Inches(0.35), PURPLE)
    add_textbox(slide, Inches(0.82), Inches(0.25), Inches(2.5), Inches(0.35),
                "DEVGTM MESH", size=11, bold=True, color=PURPLE_L)


def slide_number(slide, n, total=18):
    add_textbox(slide, Inches(11.8), Inches(7.05), Inches(1.4), Inches(0.35),
                f"{n} / {total}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def section_header(slide, label, left=Inches(0.4), top=Inches(0.75), color=PURPLE_L):
    add_textbox(slide, left, top, Inches(4), Inches(0.35),
                label.upper(), size=11, bold=True, color=color)


# ──────────────────────────────────────────────────────────
# SLIDE 01 — Cover
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
add_textbox(sl, Inches(1.5), Inches(1.6), Inches(10.3), Inches(0.5),
            "AI INNOVATION SERIES", size=14, bold=True, color=PURPLE_L, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1.5), Inches(2.1), Inches(10.3), Inches(1.6),
            "AI Use Cases in DevGTM Mesh", size=54, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(2), Inches(3.8), Inches(9.3), Inches(0.7),
            "10 AI capabilities that turn 6 integrations into one intelligent GTM system",
            size=20, color=MUTED, align=PP_ALIGN.CENTER)
# Category pills row
cats = [("Intelligence", PURPLE_L), ("Conversation", TEAL), ("Prioritization", AMBER), ("Data Foundation", BLUE_L)]
for i, (cat, col) in enumerate(cats):
    x = Inches(1.5 + i * 2.6)
    pill = add_rect(sl, x, Inches(4.8), Inches(2.4), Inches(0.5), CARD, col)
    add_textbox(sl, x, Inches(4.82), Inches(2.4), Inches(0.46), cat,
                size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
slide_number(sl, 1)

# ──────────────────────────────────────────────────────────
# SLIDE 02 — Context
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_header(sl, "Context")
add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.65),
            "Six integrations, one AI-unified platform", size=34, bold=True, color=TEXT)
add_textbox(sl, Inches(0.4), Inches(1.8), Inches(12.5), Inches(0.45),
            "DevGTM Mesh targets developer-led GTM motions. AI is the normalization and reasoning layer across all six sources.",
            size=16, color=MUTED)

integrations = [
    ("HubSpot", "CRM", "Companies, contacts, deals — baseline account record", PURPLE_L),
    ("Reo.Dev",  "Developer Intelligence", "Accounts, developers, buyers, product-usage signals", TEAL),
    ("Fireflies","Meeting Intelligence", "Transcripts, summaries, action items, sentiment", AMBER),
    ("Instantly","Email Outreach", "Campaign stats and email engagement data", BLUE_L),
    ("HeyReach", "LinkedIn Outreach", "Campaign stats and LinkedIn lead lists", RGBColor(0xf4,0x72,0xb6)),
    ("Firecrawl","Web Extraction", "Domain content extraction for research and enrichment", GREEN),
]
cols = 3
for i, (name, cat, desc, col) in enumerate(integrations):
    row, c = divmod(i, cols)
    x = Inches(0.4 + c * 4.3)
    y = Inches(2.4 + row * 2.0)
    add_rect(sl, x, y, Inches(4.1), Inches(1.75), CARD, col)
    add_textbox(sl, x + Inches(0.15), y + Inches(0.15), Inches(3.8), Inches(0.4),
                name, size=18, bold=True, color=col)
    add_textbox(sl, x + Inches(0.15), y + Inches(0.55), Inches(3.8), Inches(0.3),
                cat.upper(), size=10, bold=True, color=MUTED)
    add_textbox(sl, x + Inches(0.15), y + Inches(0.9), Inches(3.8), Inches(0.75),
                desc, size=13, color=BODY)
slide_number(sl, 2)

# ──────────────────────────────────────────────────────────
# SLIDE 03 — AI Use Case Map
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_header(sl, "AI Use Case Map")
add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.55),
            "10 use cases · 4 categories", size=30, bold=True, color=TEXT)

categories = [
    ("Intelligence", CAT_INTEL, [
        ("1", "AI Research Bot", "Claude reasoning + web context"),
        ("2", "AI Pitch Generator", "LLM personalization"),
        ("8", "Enrichment Intelligence", "LLM-driven firmographics"),
    ]),
    ("Conversation", CAT_CONV, [
        ("3", "Streaming Chat Interface", "Claude streaming API"),
        ("7", "Meeting Intelligence", "Speech-to-text + summarization"),
        ("6", "AI-First Dashboard", "Agentic UI"),
    ]),
    ("Prioritization", CAT_PRIO, [
        ("4", "Intent Scoring", "Weighted signal aggregation"),
        ("9", "SDR Auto-Assignment", "Rule + AI scoring hybrid"),
    ]),
    ("Data Foundation", CAT_DATA, [
        ("5", "Cross-Source Normalization", "Domain-match entity resolution"),
        ("10", "Web Extraction", "LLM-guided scraping"),
    ]),
]

col_w = Inches(3.15)
for ci, (cat_name, cat_col, items) in enumerate(categories):
    cx = Inches(0.3 + ci * 3.27)
    # Column header
    add_rect(sl, cx, Inches(1.75), col_w - Inches(0.1), Inches(0.55), CARD, cat_col)
    add_textbox(sl, cx, Inches(1.77), col_w - Inches(0.1), Inches(0.51),
                cat_name.upper(), size=13, bold=True, color=cat_col, align=PP_ALIGN.CENTER)
    # Items
    for ii, (num, title, sub) in enumerate(items):
        iy = Inches(2.45 + ii * 1.6)
        add_rect(sl, cx, iy, col_w - Inches(0.1), Inches(1.45), CARD, BORDER)
        add_rect(sl, cx + Inches(0.12), iy + Inches(0.15), Inches(0.45), Inches(0.45), CARD, cat_col)
        add_textbox(sl, cx + Inches(0.12), iy + Inches(0.18), Inches(0.45), Inches(0.4),
                    num, size=14, bold=True, color=cat_col, align=PP_ALIGN.CENTER)
        add_textbox(sl, cx + Inches(0.65), iy + Inches(0.12), col_w - Inches(0.82), Inches(0.5),
                    title, size=15, bold=True, color=TEXT)
        add_textbox(sl, cx + Inches(0.65), iy + Inches(0.65), col_w - Inches(0.82), Inches(0.5),
                    sub, size=12, color=MUTED)
slide_number(sl, 3)

# ──────────────────────────────────────────────────────────
# Helper: UC slide (left text / right mockup)
# ──────────────────────────────────────────────────────────
UC_SLIDES = [
    # (num, title, cat_label, cat_col, desc, bullets, mockup_title, mockup_lines, data_tags)
    (
        "01", "AI Research Bot", "Intelligence", CAT_INTEL,
        "Automates manual prospect research. Firecrawl scrapes the company domain; "
        "Claude synthesizes a structured research brief and streams it into the account panel.",
        [
            "Firecrawl extracts unstructured web content from target domain",
            "Claude reasoning synthesizes raw text into structured brief",
            "Streaming API delivers output token-by-token in real time",
            "Replaces 15–30 min of manual SDR research per account",
        ],
        "Account Research Agent — stripe.com",
        "Stripe — Payment Infrastructure (2010)\n"
        "HQ: San Francisco · 8,000+ employees\n\n"
        "Tech stack: Ruby, Go, React, AWS, Kubernetes\n"
        "Recent: Series I at $95B valuation, EU expansion\n\n"
        "Key contacts: CEO Patrick Collison, CTO David Singleton",
        ["Firecrawl scrape", "Claude claude-3-5-sonnet", "enrichment_data table"],
        4,
    ),
    (
        "02", "AI Pitch Generator", "Intelligence", CAT_INTEL,
        "Generates a personalized outreach pitch from account signals and research context. "
        "Output saved as a shareable slug-based page — contextually justified, not template-merged.",
        [
            "LLM reads normalized account data + research brief",
            "Reasons about why this account should care right now",
            "Generates pitch + saves to pitch_pages with unique slug",
            "SDR sends URL — not a generic templated email",
        ],
        "Pitch Generator — Stripe (intent score: 87)",
        "Hi Patrick — noticed Stripe's developer relations\n"
        "team grew 40% this quarter and you're expanding EU.\n\n"
        "[Product] helps teams track developer adoption\n"
        "signals before they appear in usage data.\n\n"
        "pitch.devgtm.io/stripe-q1-2025",
        ["accounts table", "enrichment_data", "signal_score", "pitch_pages table"],
        5,
    ),
    (
        "03", "Streaming Chat Interface", "Conversation", CAT_CONV,
        "A Claude-powered conversational layer over the entire GTM dataset. "
        "Real-time analyst-style Q&A via Server-Sent Events. Replaces static dashboards.",
        [
            "Claude streaming API (SSE content_block_delta events)",
            "account-ai-chat edge function injects account context as system prompt",
            "Multi-turn conversation with full GTM data as context",
            "Available globally — on account pages and dashboard",
        ],
        "DevGTM Chat — activity spike query",
        "Q: Which accounts spiked in activity this week?\n\n"
        "3 accounts with 2x+ activity spike:\n"
        "1. Datadog (+340%) — 12 dev events, 2 new buyers\n"
        "2. Vercel (+210%) — repo activity surge elevated\n"
        "3. Clerk (+190%) — meeting booked, score now 91▲",
        ["account-ai-chat edge fn", "SSE streaming", "activities table"],
        6,
    ),
    (
        "04", "Intent Scoring", "Prioritization", CAT_PRIO,
        "Weighted signal aggregation across all six integration sources. "
        "compute-intent-scores edge function writes a signal_score per account. AI explains the ranking.",
        [
            "HubSpot: deal stage, activity recency, contact engagement",
            "Reo.Dev: developer events, buyer signals, repo activity",
            "Fireflies: meeting recency, sentiment, action items",
            "Instantly / HeyReach: email opens, replies, connection rate",
        ],
        "Intent Score Breakdown — Clerk.dev",
        "Score: 91 / 100\n\n"
        "HubSpot activity  (w: 0.25)  →  88\n"
        "Reo.Dev dev sigs  (w: 0.35)  →  96\n"
        "Meeting recency   (w: 0.20)  →  90\n"
        "Outreach engage   (w: 0.20)  →  82\n\n"
        "Trend: ↑ 23 pts this week\n"
        "Reason: 4 new devs, deal active",
        ["intent_score_config", "compute-intent-scores", "signal_score field"],
        7,
    ),
    (
        "05", "Cross-Source Normalization", "Data Foundation", CAT_DATA,
        "Domain-match entity resolution creates one canonical account record from six "
        "integration sources. The prerequisite for every other AI use case in the stack.",
        [
            "HubSpot 'Stripe Inc' + Reo.Dev 'stripe.com' + Fireflies '@stripe.com' → one record",
            "domain_mappings table tracks source-to-canonical linkage",
            "Prerequisite for intent scoring, research, and pitch generation",
            "Manual override layer handles subsidiaries and rebrand cases",
        ],
        "normalize-sync-data — stripe.com",
        "HubSpot  → 'Stripe Inc'       ✓ matched\n"
        "Reo.Dev  → 'stripe.com'       ✓ matched\n"
        "Fireflies→ '@stripe.com'      ✓ matched\n"
        "Instantly→ 'stripe.com'       ✓ matched\n"
        "HeyReach → 'Stripe'           ✓ matched\n\n"
        "Canonical account: stripe.com\n"
        "5 sources merged → 1 record",
        ["accounts table", "domain_mappings", "normalize-sync-data edge fn"],
        8,
    ),
    (
        "06", "AI-First Dashboard", "Conversation", CAT_CONV,
        "Replaces the traditional metrics grid with an intent-driven interface. "
        "Global Claude chat panel with context-aware suggested prompts. Dashboard as agent.",
        [
            "Surfaces what matters — ranked by intent score, not recency",
            "Suggested prompts adapt to current account context",
            "Dashboard as agent — surfaces, suggests, executes",
            "Integration health displayed without switching tools",
        ],
        "DevGTM Dashboard — Top Accounts",
        "Top Accounts by Intent Score:\n\n"
        "1. Clerk.dev    ▪ Score: 91  ↑ HOT\n"
        "2. Datadog      ▪ Score: 87  ↑\n"
        "3. Vercel       ▪ Score: 84  →\n"
        "4. Linear       ▪ Score: 78  →\n\n"
        "Suggested: 'Why did Clerk spike?'\n"
        "Suggested: 'Draft pitch for Datadog'",
        ["signal_score ranking", "integration stats", "global chat panel"],
        9,
    ),
    (
        "07", "Meeting Intelligence", "Conversation", CAT_CONV,
        "Fireflies transcripts matched to accounts by participant email domain. "
        "Meeting summaries and sentiment scores surface on the account page.",
        [
            "Fireflies GraphQL API pulls transcript, summary, sentiment",
            "Participant email domains matched to canonical accounts table",
            "Speech-to-text + summarization creates structured meeting record",
            "SDR sees every prior call context before entering a meeting",
        ],
        "Fireflies — Clerk.dev meeting",
        "Meeting: 2025-04-22 · 34 min\n"
        "Participants: sarah@clerk.dev, alex@yourco.com\n\n"
        "Summary: Discussed auth SDK roadmap, interest in\n"
        "team plan upgrade. Sarah asked about SOC2 timeline.\n\n"
        "Sentiment: Positive (0.78)\n"
        "Action: Send SOC2 report by Friday",
        ["sync-fireflies edge fn", "Fireflies GraphQL", "domain match → accounts"],
        10,
    ),
    (
        "08", "Enrichment Intelligence", "Intelligence", CAT_INTEL,
        "LLM-driven firmographic enrichment fills the gaps that CRM records leave empty. "
        "Firecrawl + Claude infers industry, tech stack, funding stage, and social profiles.",
        [
            "Firecrawl scrapes company domain on account creation or sync",
            "Claude structured output: industry, tech stack, funding, location",
            "Live inference from company's own web presence — not static DB",
            "Results written to enrichment_data, surfaced on account page",
        ],
        "Enrichment Output — linear.app",
        "Industry: Developer Tools / PM Software\n"
        "Tech stack: TypeScript, React, Electron,\n"
        "  GraphQL, Cloudflare Workers\n"
        "Funding: Series B ($35M, 2021)\n"
        "HQ: San Francisco, CA\n"
        "GitHub: linear · 8,400 stars\n"
        "Jobs: Hiring Sr. Eng (4 roles open)",
        ["enrichment_data table", "Firecrawl extract", "Claude structured output"],
        11,
    ),
    (
        "09", "SDR Auto-Assignment", "Prioritization", CAT_PRIO,
        "Routes accounts to sales reps using a hybrid of rule-based logic and AI scoring. "
        "Human manager confirms or overrides AI recommendations.",
        [
            "Rule layer: territory, industry, employee count filters",
            "AI layer: intent score sorts priority within each rep's territory",
            "Assignment modal surfaces AI recommendations for manager review",
            "account_assignments table tracks SDR, notes, and override reason",
        ],
        "SDR Assignment Modal — Clerk.dev (score: 91)",
        "Recommended: Sarah Chen (West Coast)\n"
        "Reason: Territory match + highest capacity\n"
        "  + 3 prior Clerk touchpoints\n\n"
        "Alternatives:\n"
        "• Marcus Webb  (East) — at capacity\n"
        "• Priya Sharma (EMEA) — wrong region\n\n"
        "[ Assign to Sarah ]  [ Override ]",
        ["account_assignments", "signal_score", "rule + AI hybrid"],
        12,
    ),
    (
        "10", "Web Extraction (Firecrawl)", "Data Foundation", CAT_DATA,
        "LLM-guided scraping pulls structured data from any company website. "
        "Extraction logic expressed as a prompt — adapts automatically to site changes.",
        [
            "Firecrawl handles JS rendering and bot protection",
            "Claude interprets unstructured HTML → structured JSON fields",
            "Prompt-driven extraction adapts to site changes — no brittle selectors",
            "Feeds enrichment_data, research bot, and pitch generator",
        ],
        "Firecrawl Extraction — resend.com",
        "Product: Transactional email API for developers\n"
        "Tech: Next.js, TypeScript, React\n"
        "Open roles: 6 (3x engineering)\n"
        "Pricing: Free → Pro → Scale\n"
        "Blog: 'Why we rebuilt on Cloudflare Workers'\n"
        "Hiring signal: Strong (6 open roles)",
        ["Firecrawl API", "Claude extraction prompt", "enrichment_data write"],
        13,
    ),
]

for (num, title, cat_label, cat_col, desc, bullets, mockup_title, mockup_body, tags, slide_n) in UC_SLIDES:
    sl = add_slide(); fill_bg(sl)
    add_logo(sl)

    # Left panel
    add_textbox(sl, Inches(0.4), Inches(0.75), Inches(0.9), Inches(0.8),
                num, size=60, bold=True, color=cat_col)

    add_textbox(sl, Inches(0.4), Inches(1.55), Inches(2), Inches(0.4),
                cat_label.upper(), size=11, bold=True, color=cat_col)

    add_textbox(sl, Inches(0.4), Inches(2.0), Inches(5.8), Inches(0.85),
                title, size=28, bold=True, color=TEXT)

    add_textbox(sl, Inches(0.4), Inches(2.9), Inches(5.8), Inches(1.1),
                desc, size=15, color=BODY)

    add_textbox(sl, Inches(0.4), Inches(4.1), Inches(2), Inches(0.35),
                "AI TECHNIQUE", size=11, bold=True, color=MUTED)
    for bi, bullet in enumerate(bullets):
        add_textbox(sl, Inches(0.4), Inches(4.55 + bi * 0.55), Inches(5.8), Inches(0.5),
                    "▸  " + bullet, size=14, color=BODY)

    # Vertical divider
    add_rect(sl, Inches(6.45), Inches(0.6), Inches(0.02), Inches(6.65), BORDER)

    # Right panel — mockup window
    add_rect(sl, Inches(6.65), Inches(0.75), Inches(6.4), Inches(4.8), CARD, cat_col)
    add_textbox(sl, Inches(6.8), Inches(0.92), Inches(6.1), Inches(0.45),
                mockup_title, size=14, bold=True, color=cat_col)
    add_rect(sl, Inches(6.65), Inches(1.42), Inches(6.4), Inches(0.02), BORDER)
    add_textbox(sl, Inches(6.8), Inches(1.55), Inches(6.1), Inches(3.85),
                mockup_body, size=13, color=BODY)

    # Data tags
    add_textbox(sl, Inches(6.65), Inches(5.7), Inches(2), Inches(0.35),
                "DATA SOURCES", size=11, bold=True, color=MUTED)
    tag_x = Inches(6.65)
    for tag in tags:
        tag_w = Inches(len(tag) * 0.12 + 0.4)
        add_rect(sl, tag_x, Inches(6.1), tag_w, Inches(0.4), CARD, BORDER)
        add_textbox(sl, tag_x + Inches(0.1), Inches(6.12), tag_w - Inches(0.2), Inches(0.38),
                    tag, size=12, color=BODY)
        tag_x += tag_w + Inches(0.15)

    slide_number(sl, slide_n)

# ──────────────────────────────────────────────────────────
# SLIDE 14 — Architecture Pattern
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_header(sl, "Architecture Pattern", color=PURPLE_L)
add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.65),
            "Normalize  →  Score  →  Reason  →  Act", size=34, bold=True, color=TEXT)
add_textbox(sl, Inches(0.4), Inches(1.8), Inches(12.5), Inches(0.45),
            "Each layer depends on the one below it. You cannot reason well about data that is not normalized.",
            size=16, color=MUTED)

arch_layers = [
    ("01", "Normalize", BLUE_L, "Domain-match entity resolution — one account from six sources",
     "normalize-sync-data edge fn · domain_mappings table · prerequisite for all layers above",
     "UC 5: Cross-Source Normalization · UC 10: Web Extraction"),
    ("02", "Score",     AMBER, "Weighted signal aggregation — surfaces in-market accounts",
     "compute-intent-scores · configurable weights · signal_score written to accounts table",
     "UC 4: Intent Scoring · UC 9: SDR Auto-Assignment"),
    ("03", "Reason",    PURPLE_L, "Claude analyzes context — research, pitches, conversational Q&A",
     "Claude API (claude-3-5-sonnet) · account-ai-chat · streaming SSE · structured output",
     "UC 1: Research Bot · UC 2: Pitch Generator · UC 3: Chat Interface · UC 8: Enrichment"),
    ("04", "Act",       TEAL, "SDR assignment, outreach generation, agentic UI — the outcomes",
     "pitch_pages · account_assignments · AI-First Dashboard · automated sequences",
     "UC 6: AI Dashboard · UC 7: Meeting Intelligence"),
]

for li, (num, name, col, title, desc, use_cases) in enumerate(arch_layers):
    y = Inches(2.45 + li * 1.15)
    add_rect(sl, Inches(0.4), y, Inches(12.5), Inches(1.05), CARD, col)
    # Layer label
    add_textbox(sl, Inches(0.55), y + Inches(0.06), Inches(0.5), Inches(0.35),
                num, size=12, bold=True, color=col)
    add_textbox(sl, Inches(0.55), y + Inches(0.42), Inches(1.4), Inches(0.5),
                name, size=20, bold=True, color=col)
    # Content
    add_textbox(sl, Inches(2.1), y + Inches(0.06), Inches(10.6), Inches(0.4),
                title, size=16, bold=True, color=TEXT)
    add_textbox(sl, Inches(2.1), y + Inches(0.46), Inches(10.6), Inches(0.3),
                desc, size=12, color=MUTED)
    add_textbox(sl, Inches(2.1), y + Inches(0.74), Inches(10.6), Inches(0.28),
                use_cases, size=11, bold=True, color=col)
slide_number(sl, 14)

# ──────────────────────────────────────────────────────────
# SLIDE 15 — Why Domain Matching as Join Key
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_header(sl, "Design Decision", color=BLUE_L)
add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.65),
            "Why domain as the join key?", size=34, bold=True, color=TEXT)

# Left column
add_rect(sl, Inches(0.4), Inches(1.9), Inches(6.1), Inches(2.1), CARD, BLUE_L)
add_textbox(sl, Inches(0.55), Inches(2.0), Inches(5.8), Inches(0.45),
            "The three alternatives considered", size=16, bold=True, color=BLUE_L)
alternatives = [
    ("✗", "Company name — fails on spelling variants, subsidiaries, legal entity differences", RGBColor(0xf8,0x71,0x71)),
    ("✗", "D&B / DUNS identifier — precise but requires paid enrichment on every record", RGBColor(0xf8,0x71,0x71)),
    ("✓", "Company domain — covers vast majority of B2B accounts, in every integration, free", RGBColor(0x4a,0xde,0x80)),
]
for ai, (sym, text, col) in enumerate(alternatives):
    y = Inches(2.52 + ai * 0.5)
    add_textbox(sl, Inches(0.55), y, Inches(0.3), Inches(0.42), sym, size=16, bold=True, color=col)
    add_textbox(sl, Inches(0.9), y, Inches(5.5), Inches(0.42), text, size=13, color=BODY)

add_rect(sl, Inches(0.4), Inches(4.1), Inches(6.1), Inches(1.6), CARD, BLUE_L)
add_textbox(sl, Inches(0.55), Inches(4.2), Inches(5.8), Inches(0.45),
            "When it works well", size=16, bold=True, color=BLUE_L)
add_textbox(sl, Inches(0.55), Inches(4.7), Inches(5.8), Inches(0.9),
            "Direct company email domains — stripe.com = stripe.com across all six sources. "
            "Covers 85–90% of typical B2B account lists without manual intervention.",
            size=13, color=BODY)

# Right column
add_rect(sl, Inches(6.75), Inches(1.9), Inches(6.1), Inches(2.5), CARD, AMBER)
add_textbox(sl, Inches(6.9), Inches(2.0), Inches(5.8), Inches(0.45),
            "Known failure cases", size=16, bold=True, color=AMBER)
failures = [
    "⚠  Subsidiaries — child company email differs from parent domain in CRM",
    "⚠  Personal emails — @gmail.com / @yahoo.com contacts cannot be matched",
    "⚠  Rebrands — domain changed; old records don't reconcile automatically",
]
for fi, fail in enumerate(failures):
    add_textbox(sl, Inches(6.9), Inches(2.55 + fi * 0.55), Inches(5.8), Inches(0.48),
                fail, size=13, color=AMBER)

add_rect(sl, Inches(6.75), Inches(4.55), Inches(6.1), Inches(1.15), CARD, RGBColor(0x4a,0xde,0x80))
add_textbox(sl, Inches(6.9), Inches(4.65), Inches(5.8), Inches(0.4),
            "The mitigation", size=16, bold=True, color=RGBColor(0x4a,0xde,0x80))
add_textbox(sl, Inches(6.9), Inches(5.1), Inches(5.8), Inches(0.5),
            "domain_mappings includes a manual override layer. Operators link variant domains "
            "to the canonical account. Broad coverage now, precision via override.",
            size=13, color=BODY)
slide_number(sl, 15)

# ──────────────────────────────────────────────────────────
# SLIDE 16 — Stack Choice: Lovable + Claude
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_header(sl, "Stack Choice", color=PURPLE_L)
add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.65),
            "Deterministic plumbing + probabilistic reasoning", size=34, bold=True, color=TEXT)
add_textbox(sl, Inches(0.4), Inches(1.8), Inches(12.5), Inches(0.45),
            "Built in 26 Claude prompts inside Lovable. Clean separation between what should be deterministic and what benefits from AI flexibility.",
            size=16, color=MUTED)

# Lovable card
add_rect(sl, Inches(0.4), Inches(2.4), Inches(6.1), Inches(4.6), CARD, PURPLE_L)
add_textbox(sl, Inches(0.6), Inches(2.55), Inches(5.8), Inches(0.7),
            "Lovable", size=32, bold=True, color=PURPLE_L)
add_textbox(sl, Inches(0.6), Inches(3.28), Inches(5.8), Inches(0.35),
            "DETERMINISTIC PLUMBING", size=12, bold=True, color=MUTED)
add_textbox(sl, Inches(0.6), Inches(3.68), Inches(5.8), Inches(0.8),
            "Handles UI scaffolding, data flows, authentication, integration sync logic, "
            "and Supabase edge functions. Everything that behaves identically every time.",
            size=14, color=BODY)
lovable_attrs = [
    "Dark navy design system (hsl 232 30% 10%)",
    "26 discrete implementation steps, one prompt each",
    "Route protection + role-based access",
    "Secrets via Lovable Cloud — never hardcoded",
]
for ai, attr in enumerate(lovable_attrs):
    add_textbox(sl, Inches(0.6), Inches(4.6 + ai * 0.48), Inches(5.8), Inches(0.42),
                "▸  " + attr, size=13, color=BODY)

# Claude card
add_rect(sl, Inches(6.85), Inches(2.4), Inches(6.1), Inches(4.6), CARD, TEAL)
add_textbox(sl, Inches(7.05), Inches(2.55), Inches(5.8), Inches(0.7),
            "Claude", size=32, bold=True, color=TEAL)
add_textbox(sl, Inches(7.05), Inches(3.28), Inches(5.8), Inches(0.35),
            "PROBABILISTIC REASONING", size=12, bold=True, color=MUTED)
add_textbox(sl, Inches(7.05), Inches(3.68), Inches(5.8), Inches(0.8),
            "Handles research synthesis, pitch generation, intent explanation, and "
            "conversational Q&A. Model: claude-3-5-sonnet.",
            size=14, color=BODY)
claude_attrs = [
    "Streaming API (SSE) for real-time chat responses",
    "Structured output for enrichment fields",
    "Account context as system prompt for grounded answers",
    "ANTHROPIC_API_KEY via Lovable Cloud secrets",
]
for ai, attr in enumerate(claude_attrs):
    add_textbox(sl, Inches(7.05), Inches(4.6 + ai * 0.48), Inches(5.8), Inches(0.42),
                "▸  " + attr, size=13, color=BODY)
slide_number(sl, 16)

# ──────────────────────────────────────────────────────────
# SLIDE 17 — Takeaways
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
section_header(sl, "Key Takeaways", color=PURPLE_L)
add_textbox(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.65),
            "Three principles that make AI GTM work", size=34, bold=True, color=TEXT)

takeaways = [
    ("01", BLUE_L,   "Normalize First",
     "No amount of AI reasoning compensates for fragmented, inconsistent source data. "
     "Domain-match normalization is the prerequisite for every other use case. "
     "Build the data foundation before the AI layer."),
    ("02", PURPLE_L, "AI as the Read Layer",
     "The highest-leverage AI applications in DevGTM Mesh are read actions — "
     "research synthesis, pitch generation, conversational analytics. "
     "The AI reads the normalized data and surfaces insight. It does not replace the data layer; it interprets it."),
    ("03", TEAL,     "Agentic UI as Endpoint",
     "The goal is not a smarter dashboard — it is a system that knows what to surface, "
     "when to surface it, and can act on it when asked. "
     "The AI-first dashboard and global chat panel are the interface expression of the Normalize → Score → Reason → Act stack."),
]
for ti, (num, col, title, desc) in enumerate(takeaways):
    x = Inches(0.4 + ti * 4.3)
    add_rect(sl, x, Inches(2.0), Inches(4.1), Inches(5.0), CARD, col)
    add_textbox(sl, x + Inches(0.2), Inches(2.1), Inches(3.7), Inches(0.9),
                num, size=56, bold=True, color=col)
    add_textbox(sl, x + Inches(0.2), Inches(3.0), Inches(3.7), Inches(0.6),
                title, size=22, bold=True, color=col)
    add_textbox(sl, x + Inches(0.2), Inches(3.65), Inches(3.7), Inches(3.2),
                desc, size=14, color=BODY)
slide_number(sl, 17)

# ──────────────────────────────────────────────────────────
# SLIDE 18 — Closing / References
# ──────────────────────────────────────────────────────────
sl = add_slide(); fill_bg(sl)
add_logo(sl)
# Glows simulated with large translucent rects (pptx limitation)
add_rect(sl, -Inches(1), -Inches(1), Inches(5), Inches(5), PURPLE)

add_textbox(sl, Inches(1), Inches(1.5), Inches(11.3), Inches(0.5),
            "REFERENCES", size=14, bold=True, color=PURPLE_L, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
            "10 AI use cases. One intelligent GTM system.", size=44, bold=True,
            color=TEXT, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(2), Inches(3.55), Inches(9.3), Inches(0.65),
            "DevGTM Mesh turns six developer-focused integrations into a unified, AI-powered GTM platform. Normalize. Score. Reason. Act.",
            size=18, color=MUTED, align=PP_ALIGN.CENTER)

ref_cards = [
    ("Source Guide", PURPLE_L, "getreo.dev/resources/\ndevgtm-mesh-build-guide"),
    ("AI Model", TEAL, "Anthropic Claude\nclaude-3-5-sonnet"),
    ("Platform Stack", BLUE_L, "Lovable + Supabase\n+ 6 integrations"),
    ("Build Method", AMBER, "26 prompts in Lovable\nOne step per prompt"),
]
for ri, (label, col, val) in enumerate(ref_cards):
    x = Inches(0.7 + ri * 3.0)
    add_rect(sl, x, Inches(4.5), Inches(2.8), Inches(1.8), CARD, col)
    add_textbox(sl, x + Inches(0.18), Inches(4.62), Inches(2.4), Inches(0.4),
                label.upper(), size=11, bold=True, color=col)
    add_textbox(sl, x + Inches(0.18), Inches(5.05), Inches(2.4), Inches(0.95),
                val, size=15, color=TEXT)
slide_number(sl, 18)

# ──────────────────────────────────────────────────────────
# Save
# ──────────────────────────────────────────────────────────
out_path = "/home/shekerk/cc-best-practice-clone/presentation/devgtm-ai-use-cases/devgtm-ai-use-cases.pptx"
prs.save(out_path)
print(f"Saved {len(prs.slides)} slides to {out_path}")
